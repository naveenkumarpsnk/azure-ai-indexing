import os
import re
import logging
import json
from typing import List
from dotenv import load_dotenv
load_dotenv()

# === BLOB OPERATIONS ===
def download_blobs(connection_string, container_name, local_folder=None):
    try:
        from azure.storage.blob import BlobServiceClient
    except ImportError:
        logging.warning("azure-storage-blob not installed, skipping blob download.")
        return None

    if local_folder is None:
        local_folder = os.path.join(os.getenv("TMP", "/tmp"), "blobs")
    os.makedirs(local_folder, exist_ok=True)

    blob_service = BlobServiceClient.from_connection_string(connection_string)
    container = blob_service.get_container_client(container_name)

    for blob in container.list_blobs():
        file_path = os.path.join(local_folder, blob.name)
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        try:
            blob_client = container.get_blob_client(blob)
            with open(file_path, "wb") as f:
                f.write(blob_client.download_blob().readall())
        except Exception as e:
            logging.warning(f"Failed to download blob {blob.name}: {e}")

    print(f"Downloaded blobs to {local_folder}")
    return local_folder


# === TEXT EXTRACTION ===
def extract_text_from_pdf(path):
    try:
        import fitz
    except ImportError:
        logging.warning("PyMuPDF not installed, skipping PDF extraction.")
        return ""
    try:
        text = []
        with fitz.open(path) as doc:
            for page in doc:
                text.append(page.get_text("text"))
        return "\n".join(text)
    except Exception as e:
        logging.warning(f"Failed to extract PDF {path}: {e}")
        return ""


def extract_text_from_docx(path):
    try:
        from docx import Document
    except ImportError:
        logging.warning("python-docx not installed, skipping DOCX extraction.")
        return ""
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        logging.warning(f"Failed to extract DOCX {path}: {e}")
        return ""


def extract_text_from_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        logging.warning("python-pptx not installed, skipping PPTX extraction.")
        return ""
    try:
        text = []
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logging.warning(f"Failed to extract PPTX {path}: {e}")
        return ""


def extract_text_from_xlsx(path):
    try:
        import openpyxl
    except ImportError:
        logging.warning("openpyxl not installed, skipping XLSX extraction.")
        return ""
    try:
        text = []
        wb = openpyxl.load_workbook(path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) for cell in row if cell]
                if row_text:
                    text.append(" ".join(row_text))
        return "\n".join(text)
    except Exception as e:
        logging.warning(f"Failed to extract XLSX {path}: {e}")
        return ""


def extract_text_from_csv(path):
    try:
        import pandas as pd
    except ImportError:
        logging.warning("pandas not installed, skipping CSV extraction.")
        return ""
    try:
        df = pd.read_csv(path)
        return df.to_string()
    except Exception as e:
        logging.warning(f"Failed to extract CSV {path}: {e}")
        return ""


def extract_text(path):
    path_lower = path.lower()
    if path_lower.endswith(".pdf"):
        return extract_text_from_pdf(path)
    elif path_lower.endswith(".docx"):
        return extract_text_from_docx(path)
    elif path_lower.endswith(".pptx"):
        return extract_text_from_pptx(path)
    elif path_lower.endswith(".xlsx"):
        return extract_text_from_xlsx(path)
    elif path_lower.endswith(".csv"):
        return extract_text_from_csv(path)
    return ""


# === CHUNKING ===
def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> List[str]:
    tokens = text.split()
    chunks, start = [], 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunks.append(" ".join(tokens[start:end]))
        start += chunk_size - overlap
    return chunks


# === SEARCH INDEX ===
def ensure_index(search_endpoint, search_key, index_name):
    try:
        from azure.core.credentials import AzureKeyCredential
        from azure.search.documents.indexes import SearchIndexClient
        from azure.search.documents.indexes.models import SearchIndex, SimpleField, SearchFieldDataType, SearchableField
    except ImportError:
        logging.warning("Azure Search SDK not installed, skipping index creation.")
        return
    print("inside ensure_index")
    index_client = SearchIndexClient(search_endpoint, AzureKeyCredential(search_key))
    if index_name not in [idx.name for idx in index_client.list_indexes()]:
        fields = [
            SimpleField(name="id", type=SearchFieldDataType.String, key=True),
            SearchableField(name="text", type=SearchFieldDataType.String)
        ]
        index = SearchIndex(name=index_name, fields=fields)
        index_client.create_index(index)
        print(f"Created index: {index_name}")


# def upload_documents(search_endpoint, search_key, index_name, docs):
#     if not docs:
#         print("No documents to upload.")
#         return
#     try:
#         from azure.core.credentials import AzureKeyCredential
#         from azure.search.documents import SearchClient
#     except ImportError:
#         logging.warning("Azure Search SDK not installed, skipping document upload.")
#         return
#     print("inside upload_documents")
#     search_client = SearchClient(search_endpoint, index_name, AzureKeyCredential(search_key))
#     result = search_client.upload_documents(documents=docs)
#     print(f"Upload result: {result}")

def upload_documents(search_endpoint, search_key, index_name, docs, batch_size=500):
    if not docs:
        print("No documents to upload.")
        return

    try:
        from azure.core.credentials import AzureKeyCredential
        from azure.search.documents import SearchClient
    except ImportError:
        logging.warning("Azure Search SDK not installed, skipping document upload.")
        return

    print("inside upload_documents")

    search_client = SearchClient(search_endpoint, index_name, AzureKeyCredential(search_key))

    for i in range(0, len(docs), batch_size):
        batch = docs[i:i + batch_size]
        try:
            result = search_client.upload_documents(documents=batch)
            print(f"Uploaded batch {i // batch_size + 1} ({len(batch)} docs) → Result: {result}")
        except Exception as e:
            logging.error(f"Failed to upload batch {i // batch_size + 1}: {e}")

def sanitize_id(s: str) -> str:
# Replace invalid characters with underscore
    return re.sub(r'[^A-Za-z0-9_\-=]', "_", s)

# === MAIN FUNCTION ===
def main(): 
    print("Trigger received, starting blob processing pipeline.")

    try:
        # Environment variables
        connection_string = os.getenv("BLOB_CONNECTION_STRING")
        container_name = os.getenv("BLOB_CONTAINER_NAME")
        search_endpoint = os.getenv("SEARCH_ENDPOINT")
        search_key = os.getenv("SEARCH_KEY")
        index_name = os.getenv("SEARCH_INDEX")
        print("env variables: " +connection_string + container_name + search_endpoint + search_key + index_name)
        missing_vars = [v for v in ["BLOB_CONNECTION_STRING","BLOB_CONTAINER_NAME","SEARCH_ENDPOINT","SEARCH_KEY","SEARCH_INDEX"]
                        if not os.getenv(v)]
        if missing_vars:
            logging.warning(f"Missing environment variables: {missing_vars}")
            return f"Missing required environment variables: {missing_vars}"
            

        # Download blobs
        folder = download_blobs(connection_string, container_name)
        if not folder:
            return "Blob download failed or skipped."

        #folder = os.path.join(os.getenv("TMP", "/tmp"), "blobs")  # Assuming blobs are already downloaded for testing
        # Process files


        chunked_docs = []

        for root, _, files in os.walk(folder):
            for fname in files:
                fpath = os.path.join(root, fname)
                rel_path = os.path.relpath(fpath, folder)

                text = extract_text(fpath)
                if not text.strip():
                    continue

                chunks = chunk_text(text)
                for i, chunk in enumerate(chunks):
                    raw_id = f"{rel_path.replace(os.sep, '_')}_chunk_{i}"
                    safe_id = sanitize_id(raw_id)

                    doc = {
                        "id": safe_id,
                        "text": chunk
                    }
                    chunked_docs.append(doc)

        for doc in chunked_docs:
            print(f"ID: {doc['id']}\nTEXT:\n{doc['text'][:50]}...")

                

        # Ensure index exists and upload
        ensure_index(search_endpoint, search_key, index_name)
        upload_documents(search_endpoint, search_key, index_name, chunked_docs)

        return json.dumps({"status": "success", "docs_uploaded": len(chunked_docs)})
    

    except Exception as e:
        logging.error(f"Pipeline error: {e}", exc_info=True)
        return f"Internal Server Error: {e}"
    
if __name__ == "__main__":
    main()
